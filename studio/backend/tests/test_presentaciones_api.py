"""La presentacion de presentacion en la app: modelo, plan, deck y API.

Sin Docker ni ffmpeg: el corte real corre `studio/tools/cortar_presentacion.py` dentro
del contenedor (eso se prueba en `test_presentaciones.py`). Lo que se prueba aqui es lo
que decide si el .pptx sale bien —el plan que se le pasa al cortador, el hash
que lo caduca y la forma del deck— mas las reglas del modelo que impiden que un
proyecto acabe con slides de dos colores distintos.
"""

import time
from pathlib import Path

import pytest

from app.projects import content_hash

VALID_SCRIPT = (
    "import presentacion\n"
    "PZA = presentacion.lienzo()\n"
    "from manim import *\n"
    "class Demo(Scene):\n"
    "    def construct(self):\n"
    "        self.play(Create(Circle()))\n"
)


def _presentaciones():
    """`app.presentaciones` reimportado: el fixture `client` borra los modulos `app*`
    por cada test, y una referencia de cabecera apuntaria a la clase ANTERIOR."""
    from app import presentaciones
    return presentaciones


def _create_presentacion(authed, **overrides):
    body = {"name": "Charla de tesis", "description": "d", "quality": "ql",
            "style_block": "", "tipo": "presentacion", "fondo": "blanco"}
    body.update(overrides)
    r = authed.post("/api/projects", json=body)
    assert r.status_code == 201, r.text
    return r.json()


def _rendered_clip(authed, db, cfg, pid, title, job_id, pasos=None):
    clip = authed.post(f"/api/projects/{pid}/clips",
                       json={"title": title, "script": VALID_SCRIPT,
                             "scene": "Demo"}).json()
    media = cfg.render_jobs_dir / job_id / "media"
    (media / "videos" / "x").mkdir(parents=True, exist_ok=True)
    video = media / "videos" / "x" / "Demo.mp4"
    video.write_bytes(b"fake-mp4")
    if pasos is not None:
        (media / "pasos.json").write_text(
            '{"lienzo": {"fondo": "#ffffff"}, "pasos": %s}' % pasos)
    now = time.time()
    db.insert_job({"id": job_id, "scene": "Demo", "quality": "ql", "timeout": 120,
                   "status": "queued", "script": VALID_SCRIPT, "created_at": now,
                   "project_id": pid, "clip_id": clip["id"], "content_hash": None})
    db.update_job(job_id, status="done", started_at=now, finished_at=now,
                  video_path=str(video), size_bytes=10, resolution="854x480")
    project = db.get_project(pid)
    db.update_clip(clip["id"], job_id=job_id,
                   rendered_hash=content_hash(project["style_block"],
                                              VALID_SCRIPT, "Demo"))
    return clip


# ── el modelo ────────────────────────────────────────────────────────────────

def test_el_fondo_viaja_con_el_job_no_se_relee_del_proyecto(authed):
    """Un reintento tiene que producir el MISMO archivo que el intento
    original: si el fondo se releyera del proyecto, re-lanzar un job viejo
    daria un slide de otro color que el resto del deck."""
    from app.main import db
    pid = _create_presentacion(authed, fondo="pizarra")["id"]
    clip = authed.post(f"/api/projects/{pid}/clips",
                       json={"title": "Uno", "script": VALID_SCRIPT,
                             "scene": "Demo"}).json()
    r = authed.post(f"/api/projects/{pid}/clips/{clip['id']}/render")
    assert r.status_code == 201, r.text
    assert db.get_job(r.json()["id"])["fondo"] == "pizarra"


def test_el_fondo_se_bloquea_con_renders_vigentes(authed):
    """Cambiarlo con videos hechos dejaria un deck de dos colores."""
    from app.main import cfg, db
    pid = _create_presentacion(authed)["id"]
    assert authed.patch(f"/api/projects/{pid}",
                        json={"fondo": "marca"}).status_code == 200
    _rendered_clip(authed, db, cfg, pid, "Uno", "aaaa00001111bbbb")
    r = authed.patch(f"/api/projects/{pid}", json={"fondo": "blanco"})
    assert r.status_code == 409
    assert "el fondo" in r.json()["detail"]


def test_un_fondo_que_no_es_color_ni_nombre_es_422(authed):
    assert authed.post("/api/projects", json={
        "name": "x", "quality": "ql", "tipo": "presentacion",
        "fondo": "rm -rf /"}).status_code == 422


def test_un_color_libre_si_vale(authed):
    """El fondo lo elige quien presenta: tiene que poder dar el hex exacto de
    su plantilla, no solo elegir de una lista."""
    p = _create_presentacion(authed, fondo="#1e293b")
    assert p["fondo"] == "#1e293b"


# ── el plan ──────────────────────────────────────────────────────────────────

def test_el_plan_ordena_las_escenas_y_lista_lo_que_falta(authed):
    from app.main import cfg, db, presentacion_service
    pz = _presentaciones()

    pid = _create_presentacion(authed)["id"]
    _rendered_clip(authed, db, cfg, pid, "Uno", "aaaa00001111bbbb")
    _rendered_clip(authed, db, cfg, pid, "Dos", "cccc00002222dddd")
    authed.post(f"/api/projects/{pid}/clips",
                json={"title": "Sin render", "script": VALID_SCRIPT})

    plan = presentacion_service.plan(db.get_project(pid), pz.normaliza_opciones(None))
    assert [e["titulo"] for e in plan["escenas"]] == ["Uno", "Dos"]
    assert plan["faltan"] == ["Sin render"]
    # Rutas relativas al workspace: dentro del contenedor son /workspace/...
    assert all(not e["video"].startswith("/") for e in plan["escenas"])
    assert plan["raiz"] == "/workspace"


def test_el_plan_apunta_al_pasos_json_del_media_dir_del_job(authed):
    """`presentacion.lienzo()` lo escribe DENTRO de media_dir, que es el unico sitio
    escribible del contenedor. Si el plan apuntara a otro lado, el cortador no
    veria ningun paso y la presentacion saldria de un solo slide."""
    from app.main import cfg, db, presentacion_service
    pz = _presentaciones()
    pid = _create_presentacion(authed)["id"]
    _rendered_clip(authed, db, cfg, pid, "Uno", "aaaa00001111bbbb", pasos="[]")

    plan = presentacion_service.plan(db.get_project(pid), pz.normaliza_opciones(None))
    ruta = plan["escenas"][0]["pasos_json"]
    assert ruta.endswith("render_jobs/aaaa00001111bbbb/media/pasos.json")
    assert (cfg.workspace / ruta).is_file()


def test_el_hash_caduca_el_deck_al_re_renderizar(authed):
    """Un re-render deja la MISMA ruta y otro contenido: sin el mtime en el
    hash, el deck se quedaria diciendo «al día» con los videos viejos dentro."""
    from app.main import cfg, db, presentacion_service
    pz = _presentaciones()
    pid = _create_presentacion(authed)["id"]
    _rendered_clip(authed, db, cfg, pid, "Uno", "aaaa00001111bbbb")

    project = db.get_project(pid)
    op = pz.normaliza_opciones(None)
    antes = presentacion_service._hash_plan(presentacion_service.plan(project, op), op)

    video = Path(db.get_job("aaaa00001111bbbb")["video_path"])
    video.write_bytes(b"otro-render-distinto")
    despues = presentacion_service._hash_plan(presentacion_service.plan(project, op), op)
    assert antes != despues


def test_cambiar_de_deck_tambien_caduca(authed):
    """El GIF y el video producen archivos distintos: el mismo material con
    otra opcion no puede seguir marcando «al día»."""
    from app.main import cfg, db, presentacion_service
    pz = _presentaciones()
    pid = _create_presentacion(authed)["id"]
    _rendered_clip(authed, db, cfg, pid, "Uno", "aaaa00001111bbbb")
    project = db.get_project(pid)
    plan = presentacion_service.plan(project, pz.normaliza_opciones(None))
    con_gif = presentacion_service._hash_plan(plan, pz.normaliza_opciones({"deck": "gif"}))
    con_video = presentacion_service._hash_plan(plan, pz.normaliza_opciones({"deck": "video"}))
    assert con_gif != con_video


def test_sin_renders_no_hay_plan(authed):
    from app.main import db, presentacion_service
    pz = _presentaciones()
    pid = _create_presentacion(authed)["id"]
    authed.post(f"/api/projects/{pid}/clips",
                json={"title": "Sin render", "script": VALID_SCRIPT})
    with pytest.raises(pz.PresentacionError):
        presentacion_service.plan(db.get_project(pid), pz.normaliza_opciones(None))


# ── el deck ──────────────────────────────────────────────────────────────────

def _fragmento(tmp_path, nombre, etiqueta, dur=2.0):
    """Un fragmento como lo deja `cortar_presentacion.py`, con archivos de verdad:
    python-pptx los abre para embeberlos."""
    import struct
    import zlib
    gif = tmp_path / f"{nombre}.gif"
    # GIF 1x1 valido: python-pptx mide la imagen al insertarla.
    gif.write_bytes(bytes.fromhex(
        "47494638396101000100800000000000ffffff21f90401000000002c00000000"
        "010001000002024401003b"))
    png = tmp_path / f"{nombre}.png"
    ihdr = struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0)
    def chunk(t, d):
        return (struct.pack(">I", len(d)) + t + d
                + struct.pack(">I", zlib.crc32(t + d) & 0xFFFFFFFF))
    png.write_bytes(b"\x89PNG\r\n\x1a\n" + chunk(b"IHDR", ihdr)
                    + chunk(b"IDAT", zlib.compress(b"\x00\x00\x00\x00"))
                    + chunk(b"IEND", b""))
    return {"nombre": nombre, "etiqueta": etiqueta, "escena": "Escena",
            "duracion": dur, "gif": str(gif), "poster": str(png),
            "mp4": str(png)}


def test_el_deck_es_un_slide_por_fragmento_con_el_fondo_de_la_presentacion(tmp_path):
    from pptx import Presentation
    pz = _presentaciones()
    frs = [_fragmento(tmp_path, "01-01", "La idea"),
           _fragmento(tmp_path, "01-02", "La cifra")]
    salida = pz.construir_deck(frs, tmp_path / "deck.pptx", 16 / 9, "blanco")

    prs = Presentation(str(salida))
    assert len(prs.slides) == 2
    # 16:9 sobre 7.5" de alto: 13.333" de ancho.
    assert round(prs.slide_width / pz.EMU_POR_PULGADA, 2) == 13.33
    for slide in prs.slides:
        assert str(slide.background.fill.fore_color.rgb) == "FFFFFF"
        # El fragmento ocupa el slide ENTERO: sin franjas y, en el deck de
        # video, con un clic en cualquier parte cayendo sobre el.
        forma = slide.shapes[0]
        assert (forma.width, forma.height) == (prs.slide_width, prs.slide_height)


def test_el_nombre_del_slide_es_lo_que_el_ponente_lee(tmp_path):
    """Sale en el panel de miniaturas y en el modo presentador: es el guion."""
    from pptx import Presentation
    pz = _presentaciones()
    salida = pz.construir_deck([_fragmento(tmp_path, "01-01", "La idea")],
                               tmp_path / "d.pptx", 16 / 9, "marca")
    prs = Presentation(str(salida))
    nombre = prs.slides[0]._element.find(pz.NS_P + "cSld").get("name")
    assert "La idea" in nombre and "Escena" in nombre


def test_el_deck_de_video_pide_reproduccion_automatica(tmp_path):
    """python-pptx deja el video en «al hacer clic» (<p:cond delay="indefinite">
    dentro del nodo de medios). Sin sustituir ese arbol, cada slide se quedaria
    congelado esperando un clic que el ponente no sabe que tiene que dar."""
    from lxml import etree
    from pptx import Presentation
    pz = _presentaciones()
    salida = pz.construir_deck([_fragmento(tmp_path, "01-01", "La idea")],
                               tmp_path / "v.pptx", 4 / 3, "marca", deck="video")
    prs = Presentation(str(salida))
    timing = prs.slides[0]._element.findall(pz.NS_P + "timing")
    assert len(timing) == 1, "tiene que quedar UN solo arbol de timing"
    xml = etree.tostring(timing[0]).decode()
    assert "playFrom(0.0)" in xml
    assert 'nodeType="mainSeq"' in xml
    # 4:3 sobre 7.5" de alto: 10" de ancho.
    assert round(prs.slide_width / pz.EMU_POR_PULGADA, 2) == 10.0


def test_un_deck_sin_fragmentos_no_se_arma(tmp_path):
    pz = _presentaciones()
    with pytest.raises(pz.PresentacionError):
        pz.construir_deck([], tmp_path / "x.pptx", 16 / 9, "marca")


# ── la API ───────────────────────────────────────────────────────────────────

def test_estado_sin_material(authed):
    pid = _create_presentacion(authed)["id"]
    e = authed.get(f"/api/projects/{pid}/presentacion").json()
    assert e["estado"] == "sin_clips"
    assert e["opciones"]["deck"] == "gif"


def test_armar_sin_renders_devuelve_409(authed):
    pid = _create_presentacion(authed)["id"]
    authed.post(f"/api/projects/{pid}/clips",
                json={"title": "Sin render", "script": VALID_SCRIPT})
    assert authed.post(f"/api/projects/{pid}/presentacion", json={}).status_code == 409


def test_un_curso_no_es_una_presentacion(authed):
    """El deck vive de los pasos que anota `presentacion.paso()`: ofrecerlo en un
    curso daria un PowerPoint de clips enteros que nadie pidio."""
    from app.main import db
    pid = _create_presentacion(authed, tipo="curso", fondo="marca")["id"]
    assert db.get_project(pid)["tipo"] == "curso"
    assert authed.get(f"/api/projects/{pid}/presentacion").status_code == 409


def test_el_deck_falta_hasta_que_se_arma(authed):
    pid = _create_presentacion(authed)["id"]
    assert authed.get(f"/api/projects/{pid}/presentacion/deck").status_code == 404


def test_deck_armado_se_sirve_y_se_borra(authed):
    from app.main import cfg, db, presentacion_service
    pid = _create_presentacion(authed)["id"]
    project = db.get_project(pid)
    destino = presentacion_service.destino(project)
    destino.mkdir(parents=True, exist_ok=True)
    (destino / "deck.pptx").write_bytes(b"PK\x03\x04fake")
    (destino / "presentacion.json").write_text('{"total": 2}')

    r = authed.get(f"/api/projects/{pid}/presentacion/deck")
    assert r.status_code == 200
    assert "presentationml" in r.headers["content-type"]
    assert "charla-de-tesis.pptx" in r.headers["content-disposition"]

    assert authed.delete(f"/api/projects/{pid}/presentacion").status_code == 200
    assert not (destino / "deck.pptx").exists()
    assert authed.get(f"/api/projects/{pid}/presentacion/deck").status_code == 404


def test_presentacion_requiere_sesion(client):
    assert client.get("/api/projects/x/presentacion").status_code == 401


# ── las dos trampas que solo aparecen con la app entera montada ──────────────

def test_el_informe_del_cortador_habla_en_rutas_relativas(authed, tmp_path):
    """El cortador corre DENTRO del contenedor, donde la raiz es /workspace.
    Si el informe trajera rutas absolutas, el backend —que esta fuera— las
    abriria contra su propio sistema de archivos y no existirian.
    """
    from app.main import presentacion_service
    fr = {"nombre": "01-01", "gif": "exports/presentaciones/x/fragmentos/01-01.gif",
          "mp4": "exports/presentaciones/x/fragmentos/01-01.mp4",
          "poster": "exports/presentaciones/x/posters/01-01.png"}
    abs_ = presentacion_service._absolutas(fr)
    ws = str(presentacion_service.cfg.workspace)
    assert all(abs_[k].startswith(ws) for k in ("gif", "mp4", "poster"))
    assert abs_["nombre"] == "01-01"   # lo que no es ruta no se toca


def test_un_render_jobs_enlazado_a_otro_disco_sigue_dentro_del_workspace(
        authed, tmp_path, monkeypatch):
    """`render_jobs/` y `exports/` pueden ser enlaces a otro disco (ver
    studio/docs/ARTEFACTOS-LOCALES.md). Resolver el enlace manda la ruta fuera
    del workspace y `relative_to` la rechazaria, aunque el contenedor la vea
    perfectamente montada.
    """
    from app.main import presentacion_service
    ws = presentacion_service.cfg.workspace
    fuera = tmp_path.parent / "otro-disco" / "render_jobs"
    fuera.mkdir(parents=True, exist_ok=True)
    enlace = ws / "render_jobs_enlazado"
    if enlace.is_symlink():
        enlace.unlink()
    enlace.symlink_to(fuera)

    assert presentacion_service._rel(enlace / "abc" / "v.mp4") == \
        "render_jobs_enlazado/abc/v.mp4"


def test_un_armado_que_falla_sigue_diciendo_por_que(authed):
    """El error vivia solo mientras la tarea corria: cuando terminaba, el
    estado volvia a «sin armar» sin explicacion y no habia forma de saber que
    habia pasado."""
    from app.main import db, presentacion_service
    pid = _create_presentacion(authed)["id"]
    presentacion_service._run = {"project_id": pid, "estado": "error",
                          "error": "el corte salio con codigo 1"}
    presentacion_service._task = None            # la tarea ya termino

    e = presentacion_service.estado(db.get_project(pid))
    assert e["estado"] != "armando"       # no se queda colgado en «armando»
    assert e["run"]["error"] == "el corte salio con codigo 1"


def test_una_corrida_que_salio_BIEN_no_se_queda_pegada(authed):
    """Lo que cuenta tras un armado bueno es el informe, no el rastro de la
    corrida: dejarlo visible ensuciaria el panel para siempre."""
    from app.main import db, presentacion_service
    pid = _create_presentacion(authed)["id"]
    presentacion_service._run = {"project_id": pid, "estado": "listo"}
    presentacion_service._task = None
    assert presentacion_service.estado(db.get_project(pid))["run"] is None


# ── abrir una animacion de la Biblioteca como presentacion ──────────────────

ANIMACION = ("import sys\n"
             "sys.path.insert(0, '/workspace/studio/content/manim_extensions')\n"
             "from manim import *\n"
             "class Diagrama(Scene):\n"
             "    def construct(self):\n"
             "        self.play(Create(Circle()))\n")


def _espiar_branding(monkeypatch):
    """Anota con que `tipo` se compone cada script.

    Se espia la llamada en vez de leer el scene.py del disco: el worker de la
    cola coge el job, el render falla (en pruebas no hay runner) y el
    directorio del job se borra. Leerlo seria una carrera.
    """
    from app import jobs as app_jobs
    visto = []
    original = app_jobs.branding.aplicar

    def espia(script, tipo="curso"):
        visto.append(tipo)
        return original(script, tipo)

    monkeypatch.setattr(app_jobs.branding, "aplicar", espia)
    return visto


def _render_de(authed, pid, script=ANIMACION, scene="Diagrama"):
    clip = authed.post(f"/api/projects/{pid}/clips",
                       json={"title": "Diagrama", "script": script,
                             "scene": scene}).json()
    r = authed.post(f"/api/projects/{pid}/clips/{clip['id']}/render")
    assert r.status_code == 201, r.text
    return r.json()


def test_el_render_de_una_presentacion_se_compone_con_su_lienzo(authed, monkeypatch):
    """Es TODO el flujo de «abrir como presentación»: un proyecto de tipo
    presentacion con estilo VACIO y el script de la animacion tal cual. Si el
    tipo no llegara a branding, el 4:3 y el fondo elegidos no llegarian a manim
    y el render saldria 16:9 sobre negro sin avisar de nada.
    """
    from app.main import db
    visto = _espiar_branding(monkeypatch)
    pid = _create_presentacion(authed, style_block="", formato="clasico",
                               fondo="blanco")["id"]
    job = _render_de(authed, pid)

    assert visto == ["presentacion"]
    # Y el lienzo pedido viaja con el job, para que el runner lo pase por entorno.
    guardado = db.get_job(job["id"])
    assert (guardado["formato"], guardado["fondo"]) == ("clasico", "blanco")


def test_un_curso_con_el_mismo_script_se_compone_con_la_marca(authed, monkeypatch):
    """El tipo del proyecto es lo unico que los separa."""
    visto = _espiar_branding(monkeypatch)
    pid = _create_presentacion(authed, tipo="curso", fondo="marca",
                               style_block="")["id"]
    _render_de(authed, pid)
    assert visto == ["curso"]


def test_el_reintento_conserva_el_tipo(authed, monkeypatch):
    """Un reintento tiene que producir el MISMO archivo. El tipo se relee del
    proyecto (es inmutable), no del job."""
    pid = _create_presentacion(authed, style_block="")["id"]
    job = _render_de(authed, pid)
    visto = _espiar_branding(monkeypatch)      # solo el reintento
    reintento = authed.post(f"/api/jobs/{job['id']}/retry")
    assert reintento.status_code == 201, reintento.text
    assert visto == ["presentacion"]


def test_el_lienzo_que_de_verdad_se_compone_es_el_del_proyecto(authed):
    """Comprobacion sin espia, sobre el texto: una animacion de curso sale con
    `adaptar_escenas` y sin el bloque de marca."""
    from app import branding
    salida = branding.aplicar(ANIMACION, tipo="presentacion")
    assert "adaptar_escenas" in salida and branding.MARCADOR not in salida
