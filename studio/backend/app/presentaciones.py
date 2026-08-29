"""Presentaciones de presentacion: de los renders de un proyecto a un .pptx.

Una presentacion es una animacion para una charla de divulgacion o una defensa de
tesis. No es un clip de curso (no la narra el TTS: la narra una persona en
vivo) ni un promo (no va en bucle dentro de una app). Lo que la define es que
**avanza cuando el ponente hace clic**, y que se entrega como PowerPoint.

Como se consigue eso, que es la decision de diseno de todo el modulo:

    Cada escena se renderiza UNA sola vez, por la cola normal de jobs.
    `presentacion.paso()` anota el instante de cada punto de clic en pasos.json, y
    `cortar_presentacion.py` corta ese unico mp4 en esos instantes, dentro del
    contenedor. Cada fragmento va en SU PROPIO slide, con autoplay.

    Asi el ultimo fotograma de un fragmento ES el primero del siguiente
    (medido: 0.22/255 de diferencia media, ruido del codec), el poster del
    slide siguiente es esa misma imagen, y el empalme no se ve.

El trabajo esta partido en dos a proposito: el contenedor corta (es el unico
sitio con ffmpeg) y el backend arma el .pptx (es el unico con python-pptx).
Meter python-pptx en la imagen de manim obligaria a reconstruirla en el VPS.

Un proyecto de tipo 'presentacion' puede tener VARIAS escenas —una charla suele
llevar cinco o seis animaciones— y todas caen en el mismo deck, en el orden
de los clips.

Estados:
    sin_clips        el proyecto no tiene clips
    faltan_renders   ninguna escena tiene render vigente
    sin_armar        hay material, no hay deck
    desactualizado   cambio un render o una opcion
    al_dia           el deck corresponde a lo que hay ahora
    armando          hay una corrida en curso

OJO con el vocabulario: `pelicula.py` llama "presentaciones" a los SEGMENTOS de una
pelicula montada, una acepcion interna suya anterior a este modulo. Aqui una
presentacion es un tipo de proyecto, y sus partes se llaman fragmentos.
"""

from __future__ import annotations

import asyncio
import hashlib
import json
import time
from pathlib import Path

from .narracion import slugify
from .projects import FONDOS, FONDO_DEFECTO, specs
from .rutas import relativa_al_workspace

NOMBRE_PLAN = "plan.json"
NOMBRE_INFORME = "presentacion.json"
NOMBRE_DECK = "deck.pptx"

DECKS = ("gif", "video")
DECK_DEFECTO = "gif"

# Alto de slide fijo en 7.5", que es el de PowerPoint desde siempre; el ancho
# sale de la proporcion MEDIDA del video. Asi el fragmento ocupa el slide
# entero: sin franjas que pintar y, en el deck de video, con un clic en
# cualquier parte cayendo sobre el.
ALTO_SLIDE_PULGADAS = 7.5
EMU_POR_PULGADA = 914400


class PresentacionError(Exception):
    """No se puede armar; el mensaje va tal cual al usuario."""


def normaliza_opciones(op: dict | None) -> dict:
    op = dict(op or {})
    deck = op.get("deck") or DECK_DEFECTO
    if deck not in DECKS:
        raise PresentacionError(f"deck desconocido: {deck}")
    return {"deck": deck, "bucle": bool(op.get("bucle"))}


# ── el deck ──────────────────────────────────────────────────────────────────
#
# Dos caminos, y la diferencia entre ellos es de FIABILIDAD, no de gusto:
#
#   GIF (defecto). PowerPoint arranca un GIF animado solo al entrar al slide,
#   en todas las versiones y en las dos plataformas, sin nada de XML de por
#   medio. Es el camino que no depende de la sala.
#
#   video. Pesa menos y no se queda en 256 colores, pero que arranque solo
#   depende del arbol <p:timing> del OOXML, que aqui se escribe a mano porque
#   python-pptx no lo expone. El archivo es valido; que el autoplay se respete
#   hay que verlo UNA vez en el PowerPoint de quien presenta. Si no arrancara,
#   el video ocupa el slide entero y un clic en cualquier parte lo dispara.

TIMING_AUTOPLAY = """
<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:tnLst>
    <p:par>
      <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
          <p:seq concurrent="1" nextAc="seek">
            <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
              <p:childTnLst>
                <p:par>
                  <p:cTn id="3" fill="hold">
                    <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                    <p:childTnLst>
                      <p:par>
                        <p:cTn id="4" fill="hold">
                          <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                          <p:childTnLst>
                            <p:par>
                              <p:cTn id="5" presetID="1" presetClass="mediacall"
                                     presetSubtype="0" fill="hold" nodeType="withEffect">
                                <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                                <p:childTnLst>
                                  <p:cmd type="call" cmd="playFrom(0.0)">
                                    <p:cBhvr>
                                      <p:cTn id="6" dur="{dur_ms}" fill="hold"/>
                                      <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
                                    </p:cBhvr>
                                  </p:cmd>
                                </p:childTnLst>
                              </p:cTn>
                            </p:par>
                          </p:childTnLst>
                        </p:cTn>
                      </p:par>
                    </p:childTnLst>
                  </p:cTn>
                </p:par>
              </p:childTnLst>
            </p:cTn>
            <p:prevCondLst>
              <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:prevCondLst>
            <p:nextCondLst>
              <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
            </p:nextCondLst>
          </p:seq>
          <p:video>
            <p:cMediaNode vol="80000">
              <p:cTn id="7" fill="hold" display="0">
                <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
              </p:cTn>
              <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            </p:cMediaNode>
          </p:video>
        </p:childTnLst>
      </p:cTn>
    </p:par>
  </p:tnLst>
</p:timing>
"""

NS_P = "{http://schemas.openxmlformats.org/presentationml/2006/main}"


def color_de(fondo: str) -> str:
    """El "#rrggbb" de un fondo con nombre, o el color tal cual."""
    return FONDOS.get(fondo, fondo if str(fondo).startswith("#")
                      else FONDOS[FONDO_DEFECTO])


def _autoplay(slide, shape, dur_s: float) -> None:
    """Cambia el <p:timing> del slide por el de reproduccion automatica."""
    from lxml import etree
    xml = TIMING_AUTOPLAY.format(spid=shape.shape_id,
                                 dur_ms=max(1, int(dur_s * 1000)))
    for viejo in slide._element.findall(NS_P + "timing"):
        slide._element.remove(viejo)
    slide._element.append(etree.fromstring(xml.encode()))


def construir_deck(fragmentos: list[dict], destino: Path, proporcion: float,
                   fondo_hex: str, deck: str = DECK_DEFECTO) -> Path:
    """Arma el .pptx: un slide por fragmento, en orden.

    `fragmentos` es la lista del informe de `cortar_presentacion.py`; cada uno trae
    rutas absolutas a su mp4, su gif y su poster.
    """
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Emu

    if not fragmentos:
        raise PresentacionError("no hay fragmentos que poner en el deck")

    prs = Presentation()
    prs.slide_height = Emu(int(ALTO_SLIDE_PULGADAS * EMU_POR_PULGADA))
    prs.slide_width = Emu(int(ALTO_SLIDE_PULGADAS * proporcion * EMU_POR_PULGADA))
    rgb = RGBColor.from_string(color_de(fondo_hex).lstrip("#").upper())

    for fr in fragmentos:
        slide = prs.slides.add_slide(prs.slide_layouts[6])   # en blanco
        # El fondo se pinta aunque el fragmento lo cubra entero: si el
        # proyector escala mal y deja un borde, que sea del color de la presentacion
        # y no del blanco de la plantilla.
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb
        if deck == "video":
            forma = slide.shapes.add_movie(
                fr["mp4"], 0, 0, prs.slide_width, prs.slide_height,
                poster_frame_image=fr["poster"], mime_type="video/mp4")
            _autoplay(slide, forma, fr.get("duracion") or 1.0)
        else:
            slide.shapes.add_picture(fr["gif"], 0, 0, prs.slide_width,
                                     prs.slide_height)
        # El nombre del slide es la etiqueta del paso: es lo que el ponente
        # lee en el panel de miniaturas y en el modo presentador.
        etiqueta = fr.get("etiqueta") or ""
        escena = fr.get("escena") or ""
        nombre = f"{escena} · {etiqueta}" if escena and etiqueta else (
            etiqueta or escena or fr["nombre"])
        slide._element.find(NS_P + "cSld").set("name", nombre[:120])

    destino.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(destino))
    return destino


class PresentacionService:
    def __init__(self, cfg, db, runner) -> None:
        self.cfg = cfg
        self.db = db
        self.runner = runner
        self._task: asyncio.Task | None = None
        self._run: dict | None = None

    # ── rutas ────────────────────────────────────────────────────────────────

    def destino(self, project: dict) -> Path:
        """`exports/presentaciones/<project_id>/`.

        Por el **id** y no por el slug, igual que las peliculas: es lo unico
        que el runner valida con un regex cerrado, y renombrar el proyecto no
        deja decks huerfanos.
        """
        return self.cfg.presentaciones_dir / project["id"]

    def _rel(self, path: Path) -> str:
        """Ruta relativa al workspace (ver `rutas.py`: no se resuelven los
        enlaces antes de comparar)."""
        rel = relativa_al_workspace(path, self.cfg.workspace)
        if rel is None:
            raise PresentacionError(
                f"{path} vive fuera del workspace y el contenedor no la ve")
        return rel

    def _absolutas(self, fragmento: dict) -> dict:
        """El fragmento con sus rutas resueltas contra el workspace del host."""
        salida = dict(fragmento)
        for clave in ("mp4", "gif", "poster"):
            ruta = fragmento.get(clave)
            if ruta:
                salida[clave] = str(self.cfg.workspace / ruta)
        return salida

    def deck_path(self, project: dict) -> Path | None:
        p = self.destino(project) / NOMBRE_DECK
        return p if p.is_file() else None

    def nombre_descarga(self, project: dict) -> str:
        return f"{slugify(project['name'])}.pptx"

    # ── el plan ──────────────────────────────────────────────────────────────

    def plan(self, project: dict, op: dict) -> dict:
        clips = self.db.list_clips(project["id"])
        if not clips:
            raise PresentacionError("el proyecto no tiene clips")

        escenas, faltan = [], []
        for clip in clips:
            job = self.db.get_job(clip["job_id"]) if clip.get("job_id") else None
            if not job or job.get("status") != "done" or not job.get("video_path"):
                faltan.append(clip["title"])
                continue
            # Una presentacion es MUDA: se usa el mp4 que salio de manim, nunca la
            # mezcla con audio (que es cosa de promos).
            escenas.append({
                "titulo": clip["title"],
                "video": self._rel(Path(job["video_path"])),
                # `presentacion.lienzo()` escribe pasos.json DENTRO del media_dir del
                # job, que es el unico sitio escribible del contenedor.
                "pasos_json": (f"{self._rel(self.cfg.render_jobs_dir)}/"
                               f"{job['id']}/media/pasos.json"),
            })
        if not escenas:
            raise PresentacionError("ninguna escena de la presentacion tiene render vigente")

        return {
            "proyecto": project["name"],
            "raiz": "/workspace",
            "destino": self._rel(self.destino(project)),
            "gif": True,          # el GIF sale SIEMPRE: es el deck de respaldo
            "bucle": op["bucle"],
            "escenas": escenas,
            "faltan": faltan,
        }

    def _hash_plan(self, plan: dict, op: dict) -> str:
        """Que hace viejo a un deck: el material y como se armo.

        Se hashea con el mtime de cada archivo: un re-render deja la misma
        ruta y otro contenido.
        """
        partes = [plan["proyecto"], op["deck"], str(op["bucle"])]
        for esc in plan["escenas"]:
            for clave in ("video", "pasos_json"):
                ruta = esc[clave]
                abs_ = self.cfg.workspace / ruta
                try:
                    partes.append(f"{ruta}:{abs_.stat().st_mtime_ns}")
                except OSError:
                    partes.append(f"{ruta}:?")
        return hashlib.sha256("|".join(partes).encode()).hexdigest()[:16]

    # ── estado ───────────────────────────────────────────────────────────────

    def informe(self, project: dict) -> dict | None:
        p = self.destino(project) / NOMBRE_INFORME
        try:
            return json.loads(p.read_text()) if p.is_file() else None
        except (OSError, ValueError):
            return None

    def estado(self, project: dict, opciones: dict | None = None) -> dict:
        # La corrida sigue visible DESPUES de terminar si acabo en error. Sin
        # esto, un armado que falla vuelve a "sin armar" sin decir por que:
        # el `_run` con el mensaje solo existia mientras la tarea corria, y
        # cuando la interfaz deja de sondear ya no queda nada que ensenar.
        # Una corrida que salio bien si desaparece: lo que cuenta entonces es
        # el informe.
        corriendo = self._run
        if corriendo and corriendo.get("project_id") != project["id"]:
            corriendo = None
        elif corriendo and not self.running and corriendo.get("estado") != "error":
            corriendo = None
        informe = self.informe(project)
        deck = self.deck_path(project)
        op = normaliza_opciones(opciones or (informe or {}).get("opciones"))

        try:
            plan = self.plan(project, op)
            hash_actual = self._hash_plan(plan, op)
            escenas, faltan = len(plan["escenas"]), plan["faltan"]
            problema = None
        except PresentacionError as exc:
            hash_actual, escenas, faltan = None, 0, []
            problema = str(exc)

        if self.running and corriendo:
            estado = "armando"
        elif problema and not deck:
            estado = ("sin_clips" if "no tiene clips" in problema
                      else "faltan_renders")
        elif not deck:
            estado = "sin_armar"
        elif informe and informe.get("hash") == hash_actual:
            estado = "al_dia"
        else:
            estado = "desactualizado"

        return {
            "estado": estado,
            "problema": problema,
            "opciones": op,
            "escenas": escenas,
            "faltan": faltan,
            "informe": informe,
            "run": corriendo,
            "decks": list(DECKS),
        }

    # ── armar ────────────────────────────────────────────────────────────────

    @property
    def running(self) -> bool:
        return self._task is not None and not self._task.done()

    def start(self, project: dict, opciones: dict | None = None) -> dict:
        if self.running:
            raise PresentacionError("ya hay un deck armandose")
        op = normaliza_opciones(opciones)
        plan = self.plan(project, op)

        destino = self.destino(project)
        destino.mkdir(parents=True, exist_ok=True)
        (destino / NOMBRE_PLAN).write_text(json.dumps(plan, indent=1))

        self._run = {
            "project_id": project["id"],
            "nombre": project["name"],
            "escenas": len(plan["escenas"]),
            "deck": op["deck"],
            "iniciado": time.time(),
            "estado": "armando",
        }
        self._task = asyncio.get_event_loop().create_task(
            self._correr(project, op, self._hash_plan(plan, op)))
        return {"iniciado": True, "escenas": len(plan["escenas"]),
                "faltan": plan["faltan"]}

    async def _correr(self, project: dict, op: dict, hash_plan: str) -> None:
        destino = self.destino(project)
        try:
            informe = dict(await self.runner.cortar_presentacion(project["id"]) or {})
            fragmentos = informe.get("fragmentos") or []
            if not fragmentos:
                raise PresentacionError("el corte no dejo ningun fragmento")

            # El cortador habla en rutas relativas al workspace: dentro del
            # contenedor la raiz es /workspace y aqui es otra. Se resuelven
            # ANTES de armar el deck, que es quien abre los archivos.
            fragmentos = [self._absolutas(f) for f in fragmentos]
            proporcion = _proporcion(informe.get("resolucion"), project)
            # El fondo del render manda sobre el del proyecto: si la escena
            # eligio su color en el codigo, el slide se pinta de ESE.
            fondo = informe.get("fondo") or project.get("fondo") or FONDO_DEFECTO
            # Armar el pptx es CPU del backend, no del event loop: bloquea
            # decimas de segundo por slide y con 20 slides se notaria en toda
            # la app.
            deck = await asyncio.to_thread(
                construir_deck, fragmentos, destino / NOMBRE_DECK, proporcion,
                fondo, op["deck"])

            informe["hash"] = hash_plan
            informe["opciones"] = op
            informe["armado"] = time.time()
            informe["proporcion"] = round(proporcion, 4)
            informe["fondo_usado"] = color_de(fondo)
            informe["peso_deck"] = deck.stat().st_size
            (destino / NOMBRE_INFORME).write_text(json.dumps(informe, indent=1))
            if self._run:
                self._run["estado"] = "listo"
        except Exception as exc:  # noqa: BLE001 - el error va a la UI tal cual
            if self._run:
                self._run["estado"] = "error"
                self._run["error"] = str(exc)[:400]
        finally:
            if self._run:
                self._run["terminado"] = time.time()

    def cancel(self) -> bool:
        if not self.running:
            return False
        self._task.cancel()
        if self._run:
            self._run["estado"] = "cancelado"
        return True

    def borrar(self, project: dict) -> bool:
        """Borra el deck armado (no el material del que salio)."""
        destino = self.destino(project)
        borrado = False
        for nombre in (NOMBRE_DECK, NOMBRE_INFORME, NOMBRE_PLAN):
            p = destino / nombre
            if p.is_file():
                p.unlink()
                borrado = True
        return borrado


def _proporcion(resolucion: str | None, project: dict) -> float:
    """La proporcion del slide sale de la resolucion MEDIDA del render.

    El formato pedido es solo el respaldo: si la escena no aplico el lienzo,
    el slide tiene que caber al archivo que existe, no al que se encargo.
    """
    if resolucion and "x" in str(resolucion):
        try:
            ancho, alto = (int(v) for v in str(resolucion).split("x")[:2])
            if ancho > 0 and alto > 0:
                return ancho / alto
        except ValueError:
            pass
    sp = specs(project["quality"], project.get("formato") or "horizontal")
    return sp["width"] / sp["height"]
